from flask import Flask, render_template, request, redirect, url_for, session, flash, jsonify, send_from_directory
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
import sqlite3
import os
import mimetypes
from PIL import Image
from pptx import Presentation
from docx import Document
import openpyxl
import mammoth
import io
import base64
import uuid

app = Flask(__name__)
app.secret_key = 'your_secret_key_here'  # Change this in production

DB_PATH = os.path.join(os.path.dirname(__file__), 'users.db')
UPLOAD_FOLDER = os.path.join(os.getcwd(), 'uploads')
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

TRASH_FOLDER = os.path.join(os.getcwd(), 'uploads', 'trash')
if not os.path.exists(TRASH_FOLDER):
    os.makedirs(TRASH_FOLDER)

STAR_FILE = os.path.join(os.getcwd(), 'uploads', 'starred_files.txt')
SHARE_FILE = os.path.join(os.getcwd(), 'uploads', 'shared_files.txt')
SHARE_LINKS_FILE = os.path.join(os.getcwd(), 'uploads', 'share_links.txt')

ALLOWED_EXTENSIONS = {
    'txt', 'pdf', 'png', 'jpg', 'jpeg', 'gif', 'doc', 'docx', 
    'ppt', 'pptx', 'xls', 'xlsx', 'csv'
}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

def get_db_connection():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    conn = get_db_connection()
    conn.execute('''CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT UNIQUE NOT NULL,
        email TEXT UNIQUE NOT NULL,
        password TEXT NOT NULL
    )''')
    conn.commit()
    conn.close()

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def move_to_trash(filename):
    src = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    dst = os.path.join(TRASH_FOLDER, filename)
    if os.path.exists(src):
        os.rename(src, dst)
        return True
    return False

def list_trash_files():
    files = []
    for fname in os.listdir(TRASH_FOLDER):
        fpath = os.path.join(TRASH_FOLDER, fname)
        if os.path.isfile(fpath):
            files.append({'id': fname, 'name': fname, 'size': round(os.path.getsize(fpath)/1024, 2), 'deleted_at': ''})
    return files

def restore_from_trash(filename):
    src = os.path.join(TRASH_FOLDER, filename)
    dst = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    if os.path.exists(src):
        os.rename(src, dst)
        return True
    return False

def get_starred_files():
    if not os.path.exists(STAR_FILE):
        return set()
    with open(STAR_FILE, 'r') as f:
        return set(line.strip() for line in f if line.strip())

def add_starred_file(file_id):
    starred = get_starred_files()
    starred.add(file_id)
    with open(STAR_FILE, 'w') as f:
        for fid in starred:
            f.write(fid + '\n')

def remove_starred_file(file_id):
    starred = get_starred_files()
    if file_id in starred:
        starred.remove(file_id)
        with open(STAR_FILE, 'w') as f:
            for fid in starred:
                f.write(fid + '\n')

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/signin', methods=['GET', 'POST'])
def signin():
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        conn = get_db_connection()
        user = conn.execute('SELECT * FROM users WHERE username = ?', (username,)).fetchone()
        conn.close()
        if user and check_password_hash(user['password'], password):
            session['pending_2fa_user_id'] = user['id']
            session['pending_2fa_username'] = user['username']
            # Do NOT log in user yet, redirect to 2FA
            return redirect(url_for('twofa'))
        else:
            flash('Invalid username or password', 'danger')
    return render_template('signin.html')

@app.route('/2fa', methods=['GET', 'POST'])
def twofa():
    error = None
    if 'pending_2fa_user_id' not in session:
        return redirect(url_for('signin'))
    if request.method == 'POST':
        code = request.form['code']
        if code == '123456':
            # 2FA passed, log in user
            session['user_id'] = session.pop('pending_2fa_user_id')
            session['username'] = session.pop('pending_2fa_username')
            return redirect(url_for('dashboard'))
        else:
            error = 'Invalid 2FA code.'
    return render_template('2fa.html', error=error)

@app.route('/dashboard')
def dashboard():
    if 'user_id' not in session:
        return redirect(url_for('signin'))
    return render_template('dashboard.html', username=session.get('username'))

@app.route('/signup', methods=['GET', 'POST'])
def signup():
    if request.method == 'POST':
        username = request.form['username']
        email = request.form['email']
        password = request.form['password']
        hashed_pw = generate_password_hash(password)
        try:
            conn = get_db_connection()
            conn.execute('INSERT INTO users (username, email, password) VALUES (?, ?, ?)',
                         (username, email, hashed_pw))
            conn.commit()
            conn.close()
            flash('Account created successfully! Please sign in.', 'success')
            return redirect(url_for('signin'))
        except sqlite3.IntegrityError:
            flash('Username or email already exists.', 'danger')
    return render_template('signup.html')

@app.route('/logout')
def logout():
    session.clear()
    flash('Logged out successfully.', 'info')
    return redirect(url_for('home'))

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authenticated'}), 401
    if 'file' not in request.files:
        return jsonify({'success': False, 'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'success': False, 'error': 'No selected file'}), 400
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        save_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(save_path)
        return jsonify({'success': True, 'filename': filename})
    
    return jsonify({'success': False, 'error': 'File type not allowed'}), 400

@app.route('/files')
def list_files():
    if 'user_id' not in session:
        return jsonify([])
    files = []
    for fname in os.listdir(app.config['UPLOAD_FOLDER']):
        fpath = os.path.join(app.config['UPLOAD_FOLDER'], fname)
        if os.path.isfile(fpath):
            files.append({'id': fname, 'name': fname, 'size': round(os.path.getsize(fpath)/1024, 2), 'uploaded_at': ''})
    return jsonify(files)

@app.route('/download/<file_id>')
def download_file(file_id):
    if 'user_id' not in session:
        return redirect(url_for('signin'))
    return send_from_directory(app.config['UPLOAD_FOLDER'], file_id, as_attachment=True)

@app.route('/serve_pdf/<file_id>')
def serve_pdf(file_id):
    if 'user_id' not in session:
        return 'Not authorized', 401
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], file_id)
    if not os.path.exists(file_path):
        return 'File not found', 404
    return send_from_directory(app.config['UPLOAD_FOLDER'], file_id, mimetype='application/pdf', as_attachment=False)

@app.route('/preview/<file_id>')
def preview_file(file_id):
    if 'user_id' not in session:
        return 'Not authorized', 401

    file_path = os.path.join(app.config['UPLOAD_FOLDER'], file_id)
    if not os.path.exists(file_path):
        return 'File not found', 404

    try:
        mime_type, _ = mimetypes.guess_type(file_path)
        file_ext = os.path.splitext(file_id)[1].lower()

        # Image files
        if mime_type and mime_type.startswith('image/'):
            with Image.open(file_path) as img:
                # Resize if too large
                max_size = (800, 800)
                img.thumbnail(max_size)
                img_buffer = io.BytesIO()
                img.save(img_buffer, format=img.format)
                img_base64 = base64.b64encode(img_buffer.getvalue()).decode()
                return f'<img src="data:{mime_type};base64,{img_base64}" style="max-width:100%;max-height:60vh;" />'

        # Text files
        elif (mime_type and mime_type.startswith('text/')) or file_ext in ['.txt', '.csv']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(10000)  # Limit to first 10000 characters
                return f'<pre style="max-height:60vh;overflow:auto;white-space:pre-wrap;word-wrap:break-word;">{content}</pre>'

        # PDF files - Preview in new tab (not force download)
        elif mime_type == 'application/pdf' or file_ext == '.pdf':
            pdf_url = url_for('serve_pdf', file_id=file_id)
            return f'''<div style='text-align:center;padding:36px 0;'>
                <a href="{pdf_url}" target="_blank" class="preview-action-btn main-btn">Open PDF in New Tab</a><br>
                <a href="{url_for('download_file', file_id=file_id)}" download class="preview-action-btn secondary-btn">Download PDF</a>
            </div>
            <style>
                .preview-action-btn {{
                    display: inline-block;
                    margin: 12px 0;
                    padding: 12px 32px;
                    border-radius: 6px;
                    font-size: 1.1rem;
                    font-weight: 500;
                    text-decoration: none;
                    transition: background 0.18s, color 0.18s, box-shadow 0.18s;
                    box-shadow: 0 2px 8px rgba(30,60,114,0.06);
                }}
                .main-btn {{
                    background: linear-gradient(90deg, #2a5298 0%, #1e3c72 100%);
                    color: #fff;
                    border: none;
                }}
                .main-btn:hover {{
                    background: linear-gradient(90deg, #1e3c72 0%, #2a5298 100%);
                }}
                .secondary-btn {{
                    background: #f5f7fa;
                    color: #2a5298;
                    border: 1px solid #c2d3e8;
                    margin-top: 6px;
                }}
                .secondary-btn:hover {{
                    background: #e3e9fc;
                }}
            </style>'''

        # PowerPoint files - Preview in new tab
        elif file_ext in ['.ppt', '.pptx']:
            ppt_url = url_for('serve_ppt', file_id=file_id)
            return f'''<div style='text-align:center;padding:36px 0;'>
                <a href="{ppt_url}" target="_blank" class="preview-action-btn main-btn">Open PowerPoint in New Tab</a><br>
                <a href="{url_for('download_file', file_id=file_id)}" download class="preview-action-btn secondary-btn">Download PPT</a>
            </div>
            <style>
                .preview-action-btn {{
                    display: inline-block;
                    margin: 12px 0;
                    padding: 12px 32px;
                    border-radius: 6px;
                    font-size: 1.1rem;
                    font-weight: 500;
                    text-decoration: none;
                    transition: background 0.18s, color 0.18s, box-shadow 0.18s;
                    box-shadow: 0 2px 8px rgba(30,60,114,0.06);
                }}
                .main-btn {{
                    background: linear-gradient(90deg, #2a5298 0%, #1e3c72 100%);
                    color: #fff;
                    border: none;
                }}
                .main-btn:hover {{
                    background: linear-gradient(90deg, #1e3c72 0%, #2a5298 100%);
                }}
                .secondary-btn {{
                    background: #f5f7fa;
                    color: #2a5298;
                    border: 1px solid #c2d3e8;
                    margin-top: 6px;
                }}
                .secondary-btn:hover {{
                    background: #e3e9fc;
                }}
            </style>'''

        # Word files - Preview in new tab
        elif file_ext in ['.doc', '.docx']:
            word_url = url_for('serve_word', file_id=file_id)
            return f'''<div style='text-align:center;padding:36px 0;'>
                <a href="{word_url}" target="_blank" class="preview-action-btn main-btn">Open Word Doc in New Tab</a><br>
                <a href="{url_for('download_file', file_id=file_id)}" download class="preview-action-btn secondary-btn">Download Word</a>
            </div>
            <style>
                .preview-action-btn {{
                    display: inline-block;
                    margin: 12px 0;
                    padding: 12px 32px;
                    border-radius: 6px;
                    font-size: 1.1rem;
                    font-weight: 500;
                    text-decoration: none;
                    transition: background 0.18s, color 0.18s, box-shadow 0.18s;
                    box-shadow: 0 2px 8px rgba(30,60,114,0.06);
                }}
                .main-btn {{
                    background: linear-gradient(90deg, #2a5298 0%, #1e3c72 100%);
                    color: #fff;
                    border: none;
                }}
                .main-btn:hover {{
                    background: linear-gradient(90deg, #1e3c72 0%, #2a5298 100%);
                }}
                .secondary-btn {{
                    background: #f5f7fa;
                    color: #2a5298;
                    border: 1px solid #c2d3e8;
                    margin-top: 6px;
                }}
                .secondary-btn:hover {{
                    background: #e3e9fc;
                }}
            </style>'''

        # Excel files - Preview in new tab
        elif file_ext in ['.xls', '.xlsx']:
            excel_url = url_for('serve_excel', file_id=file_id)
            return f'''<div style='text-align:center;padding:36px 0;'>
                <a href="{excel_url}" target="_blank" class="preview-action-btn main-btn">Open Excel in New Tab</a><br>
                <a href="{url_for('download_file', file_id=file_id)}" download class="preview-action-btn secondary-btn">Download Excel</a>
            </div>
            <style>
                .preview-action-btn {{
                    display: inline-block;
                    margin: 12px 0;
                    padding: 12px 32px;
                    border-radius: 6px;
                    font-size: 1.1rem;
                    font-weight: 500;
                    text-decoration: none;
                    transition: background 0.18s, color 0.18s, box-shadow 0.18s;
                    box-shadow: 0 2px 8px rgba(30,60,114,0.06);
                }}
                .main-btn {{
                    background: linear-gradient(90deg, #2a5298 0%, #1e3c72 100%);
                    color: #fff;
                    border: none;
                }}
                .main-btn:hover {{
                    background: linear-gradient(90deg, #1e3c72 0%, #2a5298 100%);
                }}
                .secondary-btn {{
                    background: #f5f7fa;
                    color: #2a5298;
                    border: 1px solid #c2d3e8;
                    margin-top: 6px;
                }}
                .secondary-btn:hover {{
                    background: #e3e9fc;
                }}
            </style>'''

        else:
            # Binary files or unsupported formats
            file_size = os.path.getsize(file_path)
            return f'''
                <div style="text-align:center;padding:20px;">
                    <h3>File Information</h3>
                    <p>Type: {mime_type or 'Unknown'}</p>
                    <p>Size: {file_size/1024:.2f} KB</p>
                    <p>This file type can be downloaded but cannot be previewed in the browser.</p>
                    <a href="/download/{file_id}" class="download-btn" style="display:inline-block;padding:10px 20px;background:#4CAF50;color:white;text-decoration:none;border-radius:5px;">Download File</a>
                </div>
            '''

    except Exception as e:
        return f'<div style="color:#d7263d;">Preview failed: {str(e)}</div>'

@app.route('/serve_ppt/<file_id>')
def serve_ppt(file_id):
    if 'user_id' not in session:
        return 'Not authorized', 401
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], file_id)
    if not os.path.exists(file_path):
        return 'File not found', 404
    prs = Presentation(file_path)
    slides_html = []
    slides_html.append("""
        <style>
            body { margin:0; background:white; }
            .slide-container { max-height: 97vh; overflow-y: auto; padding: 20px; }
            .slide { background: white; border-radius: 5px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); margin-bottom: 20px; padding: 20px; position: relative; aspect-ratio: 16/9; width: 100%; max-width: 900px; margin-left: auto; margin-right: auto; }
            .slide-title { font-size: 30px; font-weight: bold; margin-bottom: 15px; color: #2c3e50; }
            .slide-text { font-size: 20px; color: #34495e; line-height: 1.5; }
        </style>
    """)
    for idx, slide in enumerate(prs.slides):
        title = ''
        text = ''
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if not title:
                    title = shape.text.strip()
                else:
                    text += shape.text.strip() + '<br>'
        slides_html.append(f'<div class="slide"><div class="slide-title">{title}</div><div class="slide-text">{text}</div></div>')
    return '<div class="slide-container">' + ''.join(slides_html) + '</div>'

@app.route('/serve_word/<file_id>')
def serve_word(file_id):
    if 'user_id' not in session:
        return 'Not authorized', 401
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], file_id)
    if not os.path.exists(file_path):
        return 'File not found', 404
    with open(file_path, "rb") as docx_file:
        result = mammoth.convert_to_html(docx_file)
        html = result.value
        return f'<html><head><style>body{{margin:0;background:white;}}.word-content{{max-width:900px;margin:24px auto;background:white;padding:32px 24px;border-radius:8px;box-shadow:0 2px 10px rgba(0,0,0,0.06);font-size:20px;line-height:1.7;color:#222;}}</style></head><body><div class="word-content">{html}</div></body></html>'

@app.route('/serve_excel/<file_id>')
def serve_excel(file_id):
    if 'user_id' not in session:
        return 'Not authorized', 401
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], file_id)
    if not os.path.exists(file_path):
        return 'File not found', 404
    wb = openpyxl.load_workbook(file_path, data_only=True)
    sheet = wb.active
    table_html = '<html><head><style>body{margin:0;background:white;}table{border-collapse:collapse;width:auto;background:white;font-size:18px;}td,th{border:1px solid #ccc;padding:10px;}</style></head><body><div style="overflow:auto;max-width:1100px;margin:24px auto;max-height:95vh;"><table>'
    for row in sheet.iter_rows(values_only=True):
        table_html += '<tr>' + ''.join(f'<td>{cell if cell is not None else ""}</td>' for cell in row) + '</tr>'
    table_html += '</table></div></body></html>'
    return table_html

@app.route('/delete/<file_id>', methods=['POST'])
def delete_file(file_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authenticated'}), 401
    if move_to_trash(file_id):
        return jsonify({'success': True})
    return jsonify({'success': False, 'error': 'File not found'})

@app.route('/trash')
def trash():
    if 'user_id' not in session:
        return jsonify([])
    return jsonify(list_trash_files())

@app.route('/restore/<file_id>', methods=['POST'])
def restore_file(file_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authenticated'}), 401
    if restore_from_trash(file_id):
        return jsonify({'success': True})
    return jsonify({'success': False, 'error': 'File not found in trash'})

@app.route('/star/<file_id>', methods=['POST'])
def star_file(file_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authenticated'}), 401
    add_starred_file(file_id)
    return jsonify({'success': True})

@app.route('/starred')
def starred():
    if 'user_id' not in session:
        return jsonify([])
    starred_ids = get_starred_files()
    files = []
    for fname in os.listdir(app.config['UPLOAD_FOLDER']):
        fpath = os.path.join(app.config['UPLOAD_FOLDER'], fname)
        if os.path.isfile(fpath) and fname in starred_ids:
            files.append({'id': fname, 'name': fname, 'size': round(os.path.getsize(fpath)/1024, 2), 'uploaded_at': ''})
    return jsonify(files)

@app.route('/unstar/<file_id>', methods=['POST'])
def unstar_file(file_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authenticated'}), 401
    remove_starred_file(file_id)
    return jsonify({'success': True})

@app.route('/recent')
def recent():
    if 'user_id' not in session:
        return jsonify([])
    files = []
    upload_folder = app.config['UPLOAD_FOLDER']
    # Get all files with their modification times
    file_entries = []
    for fname in os.listdir(upload_folder):
        fpath = os.path.join(upload_folder, fname)
        if os.path.isfile(fpath):
            file_entries.append((fname, os.path.getmtime(fpath)))
    # Sort by modification time, descending (most recent first)
    file_entries.sort(key=lambda x: x[1], reverse=True)
    # Return the 10 most recent files
    for fname, mtime in file_entries[:10]:
        fpath = os.path.join(upload_folder, fname)
        files.append({
            'id': fname,
            'name': fname,
            'size': round(os.path.getsize(fpath)/1024, 2),
            'uploaded_at': ''
        })
    return jsonify(files)

@app.route('/share/<file_id>', methods=['POST'])
def share_file(file_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authenticated'}), 401
    data = request.get_json()
    share_with = data.get('share_with')
    if not share_with:
        return jsonify({'success': False, 'error': 'No recipient specified'}), 400
    # Save share info (append to a file or db)
    share_db = os.path.join(app.config['UPLOAD_FOLDER'], 'shared_files.txt')
    with open(share_db, 'a') as f:
        f.write(f'{file_id},{share_with}\n')
    return jsonify({'success': True})

@app.route('/shared')
def shared_with_me():
    if 'user_id' not in session:
        return jsonify([])
    user_id = session['user_id']
    # Retrieve username and email for matching
    conn = get_db_connection()
    user = conn.execute('SELECT * FROM users WHERE id = ?', (user_id,)).fetchone()
    conn.close()
    if not user:
        return jsonify([])
    username = user['username']
    email = user['email']
    share_db = os.path.join(app.config['UPLOAD_FOLDER'], 'shared_files.txt')
    shared_files = []
    if os.path.exists(share_db):
        with open(share_db, 'r') as f:
            for line in f:
                try:
                    fid, recipient = line.strip().split(',', 1)
                except ValueError:
                    continue
                if recipient == username or recipient == str(user_id) or recipient == email:
                    fpath = os.path.join(app.config['UPLOAD_FOLDER'], fid)
                    # Try to get original filename from file metadata
                    original_name = fid
                    if os.path.isfile(fpath):
                        shared_files.append({'id': fid, 'name': original_name, 'size': round(os.path.getsize(fpath)/1024, 2)})
    return jsonify(shared_files)

@app.route('/generate_share_link/<file_id>', methods=['POST'])
def generate_share_link(file_id):
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authenticated'}), 401
    data = request.get_json()
    permission = data.get('permission', 'view')
    secret_key = data.get('secret_key')
    if not secret_key or not secret_key.isdigit() or len(secret_key) != 6:
        return jsonify({'success': False, 'error': 'Invalid secret key'}), 400
    # Generate unique token
    token = str(uuid.uuid4())
    # Save mapping: token,file_id,permission,secret_key
    share_links_db = SHARE_LINKS_FILE
    with open(share_links_db, 'a') as f:
        f.write(f'{token},{file_id},{permission},{secret_key}\n')
    link = url_for('access_shared_file', token=token, _external=True)
    return jsonify({'success': True, 'link': link})

@app.route('/shared/link/<token>', methods=['GET', 'POST'])
def access_shared_file(token):
    share_links_db = SHARE_LINKS_FILE
    if not os.path.exists(share_links_db):
        return 'Invalid or expired link', 404
    file_id, permission, secret_key = None, None, None
    with open(share_links_db, 'r') as f:
        for line in f:
            try:
                t, fid, perm, skey = line.strip().split(',', 3)
                if t == token:
                    file_id, permission, secret_key = fid, perm, skey
                    break
            except ValueError:
                continue
    if not file_id:
        return 'Invalid or expired link', 404
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], file_id)
    if not os.path.isfile(file_path):
        return 'File not found', 404
    # If GET, show secret key form
    if request.method == 'GET':
        return f'''
        <html><head><title>Enter Secret Key</title></head><body style="display:flex;align-items:center;justify-content:center;height:100vh;background:#f8f9fa;">
        <form method="POST" style="background:white;padding:36px 28px;border-radius:12px;box-shadow:0 4px 24px #1a73e822;display:flex;flex-direction:column;align-items:center;" autocomplete="off" novalidate>
            <h2 style="color:#1a73e8;margin-bottom:18px;">Enter 6-digit Secret Key</h2>
            <input name="secret_key" type="text" maxlength="6" pattern="\\d{6}" inputmode="numeric" style="font-size:1.2rem;padding:10px 16px;margin-bottom:18px;border-radius:7px;border:1px solid #dadce0;width:160px;text-align:center;letter-spacing:0.2em;" required />
            <button type="submit" style="padding:8px 24px;background:#1a73e8;color:#fff;border:none;border-radius:7px;font-size:1.1rem;cursor:pointer;">Access File</button>
        </form>
        </body></html>
        '''
    # POST: check key (robustly handle all browsers)
    user_key = request.form.get('secret_key', '').strip()
    if not user_key.isdigit() or len(user_key) != 6:
        return f'<div style="color:#e74c3c;text-align:center;padding:36px;font-size:1.25rem;">Please enter a valid 6-digit code. <a href="">Try again</a></div>'
    if user_key != secret_key:
        return f'<div style="color:#e74c3c;text-align:center;padding:36px;font-size:1.25rem;">Incorrect secret key. <a href="">Try again</a></div>'
    # Serve based on permission
    if permission == 'view':
        preview_html = preview_file(file_id)
        if isinstance(preview_html, str):
            import re
            preview_html = re.sub(r'<a[^>]*download[^>]*>.*?</a>', '', preview_html, flags=re.DOTALL)
        return preview_html
    elif permission == 'view_download':
        preview_html = preview_file(file_id)
        download_url = url_for('download_file', file_id=file_id)
        if isinstance(preview_html, str):
            import re
            preview_html = re.sub(r'<a[^>]*download[^>]*>.*?</a>', '', preview_html, flags=re.DOTALL)
            preview_html += f'<div style="text-align:center;padding:24px;"><a href="{download_url}" download class="preview-action-btn main-btn">Download</a></div>'
        return preview_html
    elif permission == 'download':
        return send_from_directory(app.config['UPLOAD_FOLDER'], file_id, as_attachment=True)
    else:
        return 'Invalid permission', 400

@app.route('/user_info')
def user_info():
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authenticated'}), 401
    user_id = session['user_id']
    conn = get_db_connection()
    user = conn.execute('SELECT * FROM users WHERE id = ?', (user_id,)).fetchone()
    conn.close()
    if not user:
        return jsonify({'success': False, 'error': 'User not found'}), 404
    return jsonify({'success': True, 'username': user['username'], 'email': user['email']})

@app.route('/change_password', methods=['POST'])
def change_password():
    if 'user_id' not in session:
        return jsonify({'success': False, 'error': 'Not authenticated'}), 401
    data = request.get_json()
    current_password = data.get('current_password')
    new_password = data.get('new_password')
    if not current_password or not new_password:
        return jsonify({'success': False, 'error': 'Both fields required.'})
    user_id = session['user_id']
    conn = get_db_connection()
    user = conn.execute('SELECT * FROM users WHERE id = ?', (user_id,)).fetchone()
    if not user or not check_password_hash(user['password'], current_password):
        conn.close()
        return jsonify({'success': False, 'error': 'Current password incorrect.'})
    hashed_pw = generate_password_hash(new_password)
    conn.execute('UPDATE users SET password = ? WHERE id = ?', (hashed_pw, user_id))
    conn.commit()
    conn.close()
    return jsonify({'success': True})

@app.route('/file_metadata/<file_id>')
def file_metadata(file_id):
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], file_id)
    if not os.path.isfile(file_path):
        return jsonify({'success': False, 'error': 'File not found'}), 404
    stat = os.stat(file_path)
    # Get owner info from DB if possible
    owner = None
    conn = get_db_connection()
    # Try to get username from session
    username = None
    if 'user_id' in session:
        user = conn.execute('SELECT username FROM users WHERE id = ?', (session['user_id'],)).fetchone()
        if user:
            username = user['username']
    # Fallback to DB lookup if files table exists
    if not username and 'files' in [r[0] for r in conn.execute("SELECT name FROM sqlite_master WHERE type='table'")]:
        owner_row = conn.execute('SELECT username FROM users WHERE id = (SELECT user_id FROM files WHERE filename = ?)', (file_id,)).fetchone()
        if owner_row:
            username = owner_row['username']
    conn.close()
    import datetime
    created_dt = datetime.datetime.fromtimestamp(stat.st_ctime)
    modified_dt = datetime.datetime.fromtimestamp(stat.st_mtime)
    # Format size
    def format_size(size):
        for unit in ['bytes','KB','MB','GB','TB']:
            if size < 1024.0:
                return f"{size:.2f} {unit}" if unit != 'bytes' else f"{int(size)} bytes"
            size /= 1024.0
        return f"{size:.2f} PB"
    metadata = {
        'success': True,
        'name': file_id,
        'size': format_size(stat.st_size),
        'created': created_dt.strftime('%Y-%m-%d %H:%M:%S'),
        'modified': modified_dt.strftime('%Y-%m-%d %H:%M:%S'),
        'owner': username or 'Unknown',
        'permissions': oct(stat.st_mode)[-3:]
    }
    return jsonify(metadata)

if __name__ == '__main__':
    init_db()
    app.run(debug=True)
